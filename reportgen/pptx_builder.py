"""Сборка презентации Q1 2026 КЛОНИРОВАНИЕМ готовых слайдов шаблона
«Контур Blue 2023 Montserrat» с подменой текста.

Почему так: в шаблоне нет нативных графиков (диаграммы вставляются
картинками из Google-таблиц), а данные показываются через готовые
слайды-образцы. Чтобы выглядело как настоящая преза «в этом формате»
без самодеятельности, мы берём родные слайды шаблона и меняем в них
только текст — позиции, шрифты, цвета, декор остаются из шаблона.

Используемые образцы (1-based в шаблоне):
  12  Титул 1
  52  Отбивка / «Заголовок раздела»  (разделители)
  25  «6 важных цифр»                 (KPI-слайды)
  60  «Таблица»                       (таблицы данных)
  22  «Заголовок в 1 строку + текст»  (выводы списком)
  50  Финальный слайд_1
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

ACCENT = RGBColor(0x22, 0x91, 0xFF)
NAVY = RGBColor(0x15, 0x31, 0x77)
INK = RGBColor(0, 0, 0)
GRAY = RGBColor(0x5A, 0x65, 0x73)
GRAYBAR = RGBColor(0xC2, 0xC8, 0xD0)
PLATE = RGBColor(0xF1, 0xF1, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x26, 0xAD, 0x50)
RED = RGBColor(0x66, 0x14, 0x29)
FONT = "Montserrat"

L_BODY = "Заголовок в 1 строку + текст"   # брендовый layout для кастомных слайдов

# индексы образцов (0-based)
SRC_TITLE = 11
SRC_DIVIDER = 51
SRC_SIX = 24          # «6 важных цифр»
SRC_TWO = 23          # «2 цифры крупно»
SRC_TABLE = 59        # «Таблица»
SRC_BODY = 21         # «Заголовок в 1 строку + текст» (маркированный список)
SRC_FINAL = 49

# карта «6 важных цифр»: (idx числа, idx подписи) в порядке чтения
SIX_MAP = [(6, 5), (13, 9), (4, 3), (8, 7), (15, 14), (2, 1)]


def clone_slide(prs, src_index):
    """Создать новый слайд — точную копию слайда-образца (с картинками)."""
    src = prs.slides[src_index]
    dst = prs.slides.add_slide(src.slide_layout)
    # убрать автодобавленные плейсхолдеры — принесём элементы образца
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)  # noqa: SLF001
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))  # noqa: SLF001
    # перенести связи картинок (blip r:embed)
    for blip in dst.shapes._spTree.iter(qn('a:blip')):  # noqa: SLF001
        rid = blip.get(qn('r:embed'))
        if rid and rid in src.part.rels:
            rel = src.part.rels[rid]
            new = dst.part.relate_to(rel.target_part, rel.reltype)
            blip.set(qn('r:embed'), new)
    return dst


def _set_run_text(ph, text, *, color=None):
    """Заменить текст плейсхолдера, СОХРАНИВ форматирование шаблона."""
    tf = ph.text_frame
    lines = text.split("\n") if isinstance(text, str) else list(text)
    p0 = tf.paragraphs[0]
    # эталон форматирования — первый run
    proto = p0.runs[0] if p0.runs else None
    # очистить лишние параграфы
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    # первый параграф
    _fill_para(p0, lines[0], proto, color)
    # остальные строки — новые параграфы с тем же форматом
    for ln in lines[1:]:
        np = tf.add_paragraph()
        _fill_para(np, ln, proto, color)


def _fill_para(p, text, proto, color):
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for br in p._p.findall(qn('a:br')):   # убрать переносы строк из шаблона
        p._p.remove(br)
    r = p.add_run(); r.text = text
    if proto is not None:
        rp = proto._r.get_or_add_rPr()
        r._r.insert(0, copy.deepcopy(rp))
    r.font.name = FONT
    if color is not None:
        r.font.color.rgb = color


def _ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set(slide, idx, text, color=None):
    ph = _ph(slide, idx)
    if ph is not None:
        _set_run_text(ph, text, color=color)


def _remove_ph(slide, idx):
    ph = _ph(slide, idx)
    if ph is not None:
        ph._element.getparent().remove(ph._element)  # noqa: SLF001


# ── высокоуровневые слайды ────────────────────────────────────────

def title_slide(prs, title, subtitle, tagline):
    s = clone_slide(prs, SRC_TITLE)
    _remove_ph(s, 4)            # убрать пустой блок под фото
    _set(s, 0, title)
    _set(s, 1, subtitle)
    _set(s, 2, tagline)
    _set(s, 3, "kontur.ru")
    return s


def divider(prs, title):
    s = clone_slide(prs, SRC_DIVIDER)
    if s.shapes.title:
        _set_run_text(s.shapes.title, title)
    else:
        _set(s, 0, title)
    return s


def _ensure_title(slide, text):
    """Заголовок: заполнить title-плейсхолдер; если его нет (как в образце
    «6 важных цифр») — принести его из layout, чтобы стиль был из шаблона."""
    if _ph(slide, 0) is not None:
        _set_run_text(_ph(slide, 0), text)
        return
    lt = None
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.idx == 0:
            lt = ph; break
    if lt is not None:
        slide.shapes._spTree.append(copy.deepcopy(lt._element))  # noqa: SLF001
        _set_run_text(_ph(slide, 0), text)


def six_numbers(prs, title, items):
    """items — до 6 кортежей (число, подпись)."""
    s = clone_slide(prs, SRC_SIX)
    _ensure_title(s, title)
    for (num_idx, desc_idx), it in zip(SIX_MAP, items):
        val, desc = it
        _set(s, num_idx, val)
        _set(s, desc_idx, desc)
    # лишние ячейки (если меньше 6) — удалить плейсхолдеры (иначе торчат подсказки)
    for (num_idx, desc_idx) in SIX_MAP[len(items):]:
        _remove_ph(s, num_idx)
        _remove_ph(s, desc_idx)
    return s


def two_numbers(prs, title, left, right):
    """Слайд «2 цифры крупно»: left/right = (число, подпись)."""
    s = clone_slide(prs, SRC_TWO)
    _ensure_title(s, title)
    (ln, ld), (rn, rd) = left, right
    _set(s, 2, ln); _set(s, 1, ld)
    _set(s, 3, rn); _set(s, 4, rd)
    return s


def final_slide(prs, title):
    s = clone_slide(prs, SRC_FINAL)
    if s.shapes.title:
        _set_run_text(s.shapes.title, title)
    return s


def bullets_slide(prs, title, bullets):
    s = clone_slide(prs, SRC_BODY)
    _set(s, 0, title)
    body = _ph(s, 1)
    if body is not None:
        _set_run_text(body, list(bullets))
    return s


def table_slide(prs, title, headers, rows, *, row_h=Cm(0.78), fs=11):
    """Клон слайда «Таблица»: заголовок из шаблона, таблицу строим заново
    в брендовом стиле (синяя шапка, зебра)."""
    s = clone_slide(prs, SRC_TABLE)
    _set(s, 0, title)
    # удалить пример-таблицу и поясняющую плашку/стрелки из образца
    for sh in list(s.shapes):
        if sh.has_table or sh.shape_type == 13:   # таблица или картинка-пример
            sh._element.getparent().remove(sh._element)  # noqa: SLF001
        elif sh.has_text_frame and "выделить" in sh.text_frame.text.lower():
            sh._element.getparent().remove(sh._element)  # noqa: SLF001
    _brand_table(s, Cm(1.63), Cm(5.2), Cm(30.6), headers, rows, row_h=row_h, fs=fs)
    return s


def _brand_table(slide, x, y, w, headers, rows, *, row_h, fs):
    n = len(rows) + 1; cols = len(headers)
    t = slide.shapes.add_table(n, cols, x, y, w, row_h * n).table
    for r in t.rows:
        r.height = row_h
    tblPr = t._tbl.tblPr  # noqa: SLF001
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')

    def cell(ri, ci, val, *, head=False):
        c = t.cell(ri, ci)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Cm(0.2); c.margin_right = Cm(0.2)
        c.margin_top = Cm(0.04); c.margin_bottom = Cm(0.04)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT if head else (PLATE if ri % 2 == 0 else WHITE)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = str(val)
        r.font.name = FONT; r.font.size = Pt(fs)
        if head:
            r.font.bold = True; r.font.color.rgb = WHITE
        else:
            txt = str(val).strip()
            col = INK
            if "%" in txt and txt.startswith("+"): col = GREEN
            elif "%" in txt and (txt.startswith("−") or txt.startswith("-")): col = RED
            r.font.color.rgb = col; r.font.bold = (ci == 0)

    for ci, h in enumerate(headers):
        cell(0, ci, h, head=True)
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell(ri, ci, val)
    return t


# ── НЕТИПОВЫЕ слайды: кастом на брендовом layout'е (график) ────────

def _layout(prs, name):
    for L in prs.slide_layouts:
        if L.name == name:
            return L
    return prs.slide_layouts[0]


def _style_chart(ch, multi):
    ch.has_title = False
    ch.has_legend = multi
    if multi:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = PLATE
    va.major_gridlines.format.line.width = Pt(0.75)
    va.minimum_scale = 0
    va.format.line.fill.background()
    va.tick_labels.font.size = Pt(10); va.tick_labels.font.name = FONT; va.tick_labels.font.color.rgb = GRAY
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.name = FONT; ca.tick_labels.font.color.rgb = GRAY


def chart_slide(prs, title, kind, cats, series):
    """Нетиповой слайд на брендовом layout'е: заголовок-плейсхолдер +
    нативный график в фирменных цветах (kind='bar'|'line')."""
    s = prs.slides.add_slide(_layout(prs, L_BODY))
    _remove_ph(s, 1)                       # убрать body-подсказку
    if s.shapes.title:
        _set_run_text(s.shapes.title, title)
    cd = CategoryChartData(); cd.categories = cats
    for ser in series:
        cd.add_series(ser["name"], ser["data"])
    ctype = XL_CHART_TYPE.LINE if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = s.shapes.add_chart(ctype, Cm(1.63), Cm(5.4), Cm(30.6), Cm(11.5), cd)
    ch = gf.chart
    _style_chart(ch, len(series) > 1)
    cols = [ACCENT, NAVY, GREEN]
    if kind == "line":
        for i, ser in enumerate(ch.series):
            col = GRAYBAR if series[i].get("plan") else cols[i % len(cols)]
            ser.format.line.color.rgb = col
            ser.format.line.width = Pt(2.25 if series[i].get("plan") else 3.25)
            if series[i].get("forecast"):
                ser.format.line.dash_style = 7
            m = ser.marker; m.style = 8; m.size = 7
            m.format.fill.solid(); m.format.fill.fore_color.rgb = col
            m.format.line.color.rgb = col
    else:
        ch.plots[0].gap_width = 70
        for i, ser in enumerate(ch.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = GRAYBAR if series[i].get("plan") else cols[i % len(cols)]
            ser.format.line.fill.background()
        pl = ch.plots[0]; pl.has_data_labels = True
        pl.data_labels.number_format = '0.0'; pl.data_labels.number_format_is_linked = False
        pl.data_labels.font.size = Pt(10); pl.data_labels.font.name = FONT; pl.data_labels.font.bold = True
        pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    return s


# ── СБОРКА ────────────────────────────────────────────────────────


# ── СБОРКА (зеркалит HTML-дек) ────────────────────────────────────

def build(template_path, out_path):
    prs = Presentation(template_path)
    originals = list(prs.slides._sldIdLst)  # noqa: SLF001
    M = ["Янв","Фев","Мар","Апр","Май","Июн","Июл","Авг","Сен","Окт","Ноя","Дек"]

    title_slide(prs, "Итоги Q1 2026",
                "Квартальный отчёт по продажам · Контур",
                "Розница · Общепит · Кассы · ОФД · Маркет · прогноз 2026")

    six_numbers(prs, "Ключевые метрики квартала Q1 2026", [
        ("68,2 млн ₽", "Выручка факт"),
        ("92 %", "Выполнение плана (74,6 млн)"),
        ("6 859", "Количество оплат · 101% плана"),
        ("9,9 тыс ₽", "Средний чек"),
        ("−6,3 млн ₽", "Недобор к плану"),
        ("85 %", "Март — провал к плану"),
    ])
    table_slide(prs, "Бизнес-юнит: план и факт по месяцам, млн ₽",
                ["Месяц","План","Факт","Выполнение","Оплат"],
                [["Январь","22,33","21,12","95%","2 123"],
                 ["Февраль","24,75","23,69","96%","2 510"],
                 ["Март","27,47","23,42","85%","2 226"],
                 ["Итого Q1","74,56","68,23","92%","6 859"]])

    divider(prs, "Направления")
    table_slide(prs, "Q1 2026 по направлениям · план 74,6 → факт 68,2 млн (92%)",
                ["Направление","Факт Q1, млн","Доля","Ключевая динамика"],
                [["Розница","~32,5","48%","Оптим. Розница +9% QoQ; Премиум −42% YoY"],
                 ["Кассы","~22,2","33%","Перерег. ККТ +48% YoY; Терминал +29% QoQ"],
                 ["Общепит","~12,8","19%","Оптим. Общепит +44% QoQ — растёт"],
                 ["ОФД (срез)","11,7","—","+4,2% QoQ, но −4,9% YoY"]])
    table_slide(prs, "Розница · ~32,5 млн ₽ (48% БЮ)",
                ["Тариф","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["Оптимальный Розница","7,13","6,13","6,66","+9%","−7%"],
                 ["Все госсистемы","3,97","3,80","3,47","−9%","−13%"],
                 ["Базовый Розница","1,89","1,59","1,41","−11%","−25%"],
                 ["Премиум Розница","1,64","1,13","0,95","−16%","−42%"],
                 ["Маркировка","1,39","3,26","1,86","−43%","+34%"]])
    table_slide(prs, "Общепит · ~12,8 млн ₽ (19% БЮ) · растёт",
                ["Тариф","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["Оптимальный Общепит","1,20","0,89","1,29","+44%","+8%"],
                 ["КМ модификатор ЕГАИС","1,17","0,99","1,09","+10%","−7%"],
                 ["Базовый Общепит","—","—","0,49","·","·"]])
    table_slide(prs, "Кассы · ~22,2 млн ₽ (33% БЮ) · самое здоровое",
                ["Тариф","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["Терминал","2,63","2,14","2,75","+29%","+5%"],
                 ["Модификатор Маркировка","2,46","2,30","2,55","+11%","+4%"],
                 ["Лицензия Атол ИТС","1,23","1,39","1,59","+14%","+29%"],
                 ["Перерегистрация ККТ","1,05","1,30","1,56","+20%","+48%"],
                 ["Атол ИТС","1,48","1,00","1,24","+24%","−17%"],
                 ["Сканер","1,21","1,37","1,16","−15%","−5%"]])
    table_slide(prs, "ОФД · 11,69 млн ₽ (+4,2% QoQ · −4,9% YoY) · 2 122 оплаты",
                ["Тариф","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["ОФД-36","4,77","4,79","3,93","−18%","−18%"],
                 ["ОФД-15","4,17","3,00","3,73","+24%","−11%"],
                 ["ОФД-13","2,72","2,73","3,38","+24%","+24%"],
                 ["ОФД-24 и прочие","—","—","0,65","·","·"]])

    table_slide(prs, "Топ тарифов Маркета: QoQ и YoY",
                ["Тариф (Контур.Маркет)","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["Оптимальный Розница","7,13","6,13","6,66","+9%","−7%"],
                 ["Все госсистемы","3,97","3,80","3,47","−9%","−13%"],
                 ["Настройка Контур.Маркет","1,59","1,45","2,93","+102%","+85%"],
                 ["Маркировка","1,39","3,26","1,86","−43%","+34%"],
                 ["Базовый Розница","1,89","1,59","1,41","−11%","−25%"],
                 ["Оптимальный Общепит","1,20","0,89","1,29","+44%","+8%"],
                 ["КМ модификатор ЕГАИС","1,17","0,99","1,09","+10%","−7%"],
                 ["Премиум Розница","1,64","1,13","0,95","−16%","−42%"],
                 ["Базовый Услуги","1,04","0,67","0,86","+29%","−18%"]], row_h=Cm(0.78), fs=10)
    table_slide(prs, "Топ тарифов ОФД: QoQ и YoY",
                ["Тариф (Контур.ОФД)","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["ОФД-36","4,77","4,79","3,93","−18%","−18%"],
                 ["ОФД-15","4,17","3,00","3,73","+24%","−11%"],
                 ["ОФД-13","2,72","2,73","3,38","+24%","+24%"],
                 ["ОФД-24","—","—","0,46","·","·"],
                 ["Код активации ОФД-13","—","—","0,10","·","·"],
                 ["ОФД-18","—","—","0,07","·","·"],
                 ["Контур.ОФД-7 дней","—","—","0,02","·","·"]])
    table_slide(prs, "Топ тарифов Касс: QoQ и YoY",
                ["Тариф (Кассы)","Q1 2025","Q4 2025","Q1 2026","QoQ","YoY"],
                [["Терминал","2,63","2,14","2,75","+29%","+5%"],
                 ["Модификатор Маркировка","2,46","2,30","2,55","+11%","+4%"],
                 ["Лицензия Атол ИТС","1,23","1,39","1,59","+14%","+29%"],
                 ["Перерегистрация ККТ","1,05","1,30","1,56","+20%","+48%"],
                 ["Атол ИТС","1,48","1,00","1,24","+24%","−17%"],
                 ["Сертификат","0,86","0,66","1,23","+88%","+43%"],
                 ["Сканер","1,21","1,37","1,16","−15%","−5%"]])

    bullets_slide(prs, "Факторы квартала: успехи и сложности", [
        "Успех — маркировочные волны: «Маркировка» +34% YoY, спрос толкает продажи.",
        "Успех — кросс-продуктовые запуски: «Настройка Контур.Маркет» +102% QoQ / +85% YoY.",
        "Успех — Кассы: Перерег. ККТ +48% YoY, Терминал +29% QoQ.",
        "Успех — «тс ПИОТ»: [добавить суть инициативы и вклад].",
        "Сложность — март 85% плана: план вырос до 27,5 млн, факт на уровне февраля.",
        "Сложность — Маркет −21% YoY; Премиум/Базовый Розница −42%/−25% YoY.",
        "Сложность — ОФД −4,9% YoY; реклама ROMI ≈ −50%, CPO Маркета 46,8 тыс ₽.",
    ])

    divider(prs, "Прогноз 2026 · Маркет и ОФД")
    two_numbers(prs, "Маркет — выполнение плана 2026",
                ("119,0 млн ₽", "Годовой план (Продукт = Маркет). Факт за 5 мес — 44,5 млн (91% плана периода), −21% YoY."),
                ("91 %", "Прогноз года: 108,3 млн ₽ (база = темп YTD). Коридор 86–96%."))
    chart_slide(prs, "Маркет — план / факт / прогноз 2026, млн ₽", "line", M,
        [{"name":"План","plan":True,"data":[9.78,9.46,10.98,10.45,8.36,10.32,10.13,9.01,8.79,9.95,9.71,12.10]},
         {"name":"Факт","data":[8.24,9.20,10.26,8.99,7.78,None,None,None,None,None,None,None]},
         {"name":"Прогноз","forecast":True,"data":[None,None,None,None,7.78,9.41,9.24,8.22,8.02,9.08,8.86,11.04]}])
    two_numbers(prs, "ОФД — выполнение плана 2026",
                ("37,5 млн ₽", "Годовой план (проект п453). Факт за 5 мес — 18,9 млн (123% плана периода), −16% YoY."),
                ("115 %", "Прогноз года: 43,3 млн ₽ (база). Коридор 103–124%."))
    chart_slide(prs, "ОФД — план / факт / прогноз 2026, млн ₽", "line", M,
        [{"name":"План","plan":True,"data":[2.81,3.11,3.18,3.24,2.97,3.04,3.03,2.99,2.90,3.28,3.20,3.79]},
         {"name":"Факт","data":[4.05,4.12,4.18,3.71,2.84,None,None,None,None,None,None,None]},
         {"name":"Прогноз","forecast":True,"data":[None,None,None,None,2.84,3.6,3.4,3.1,3.1,3.4,3.2,4.5]}])
    six_numbers(prs, "Прогноз по лидам 2026 (мой расчёт)", [
        ("16,1 тыс", "Маркет — лиды (нужно ~17,8 тыс)"),
        ("36 %", "Маркет — CR2 лид→оплата"),
        ("5,8 тыс", "Маркет — оплаты прогноз"),
        ("15,4 тыс", "ОФД — лиды (минимум ~15,0 тыс)"),
        ("46 %", "ОФД — CR2 лид→оплата"),
        ("7,1 тыс", "ОФД — оплаты прогноз"),
    ])

    divider(prs, "Контекстная реклама")
    table_slide(prs, "Контекст: Маркет vs ОФД (янв 2025 – май 2026)",
                ["Показатель","Маркет","ОФД","Разница"],
                [["Расход","14,7 млн ₽","7,1 млн ₽","Маркет ×2,1"],
                 ["Оплаты привлечено","315","397","ОФД +26%"],
                 ["CPO — стоимость оплаты","46,8 тыс ₽","17,9 тыс ₽","ОФД ×2,6 дешевле"],
                 ["CPL — средний","10,0 тыс ₽","4,5 тыс ₽","ОФД ×2,2 дешевле"],
                 ["Выручка по целевым лидам","3,9 млн ₽","2,8 млн ₽","—"],
                 ["ДРР (расход / выручка)","≈ 3,8×","≈ 2,5×","ОФД эффективнее"]],
                row_h=Cm(1.0), fs=12)
    table_slide(prs, "Контекст: привлечение по продуктам (оплаты по месяцам)",
                ["Месяц","ОФД","Маркет","Бандл"],
                [["М1","54","63","52"],["М2","65","95","46"],["М3","50","80","60"],
                 ["М4","52","36","43"],["М5","37","28","68"],["М6","8","13","1"],
                 ["Итого","266","315","270"]])

    divider(prs, "Маркетинг: каналы и планы 2026")
    table_slide(prs, "Контекст: годовая динамика 2023–2026 (РК Маркет + бандл)",
                ["Показатель","2023","2024","2025","2026 план"],
                [["Расход, млн ₽","32,6","32,6","24,7","17,0"],
                 ["Лиды","3 986","3 572","3 205","2 900"],
                 ["Оплаты","536","527","538","500"],
                 ["CR в оплату","13,4%","14,8%","16,8%","17,2%"],
                 ["Выручка, млн ₽","18,0","18,8","19,7","19,0"],
                 ["Выручка − расход, млн ₽","−14,5","−13,7","−5,0","+2,0"]])
    six_numbers(prs, "Брендформанс и SEO — планы 2026", [
        (">2 млн ₽", "Брендформанс — факт окт–дек 2025"),
        ("~11 млн ₽", "Брендформанс — план 2026 (6,5 — Маркет)"),
        ("+3 %", "SEO 2025 (при норме +20%)"),
        ("~+10 млн ₽", "SEO — план 2026"),
    ])
    six_numbers(prs, "Аудит платной рекламы (М. Петрова)", [
        ("−12 %", "выручки при отключении всей рекламы (~12 млн ₽)"),
        ("32 %", "выручки ОФД даёт реклама на Бандле"),
        ("≈ −50 %", "ROMI реалистично"),
        ("+15 %", "ROMI в мае (оптимизация + продажи)"),
    ])

    bullets_slide(prs, "Итоги Q1 2026: ключевые выводы", [
        "Q1: 92% плана. Январь–февраль 95–96%, март провален (85%).",
        "Маркет: −21% YoY, база прогноза 2026 ≈ 91% плана (темп YTD сохраняется).",
        "ОФД: факт YTD 123% плана периода, прогноз 2026 ≈ 115% — перевыполнение.",
        "Лиды (мой расчёт): Маркет ~16,1 тыс (нужно 17,8), ОФД ~15,4 тыс (минимум 15,0).",
        "Реклама: ОФД-контекст в 2,6× эффективнее Маркета по стоимости оплаты.",
    ])
    final_slide(prs, "Благодарю\nза внимание!")

    for sid in originals:
        prs.part.drop_rel(sid.get(qn('r:id')))
        prs.slides._sldIdLst.remove(sid)  # noqa: SLF001
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


if __name__ == "__main__":
    import sys
    t = sys.argv[1] if len(sys.argv) > 1 else "Шаблон презентации Контур Blue 2023_16x9_Montserrat (2).pptx"
    o = sys.argv[2] if len(sys.argv) > 2 else "Итоги Q1 2026 — Контур.pptx"
    build(t, o)
    print("OK ->", o)
