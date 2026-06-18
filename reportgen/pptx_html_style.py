"""PPTX вариант B — кастомная вёрстка «как в HTML-деке».

Не клонирует слайды-образцы шаблона: каждый слайд собирается с нуля на
чистом макете, чтобы визуально совпадать с HTML — карточки KPI с серыми
плашками и синим/красным акцентом дельты, синие шапки таблиц, графики
прогноза с осями и легендой справа, отбивки с тёмно-синим фоном.

Цвета — из брендовой темы (theme1.xml): accent1 #2291FF, accent2 #153177.
Шрифт — Montserrat (наследуется из шаблона), все шейпы явно задают FONT.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# Палитра (как в HTML, бренд Контур)
ACCENT = RGBColor(0x22, 0x91, 0xFF)
ACCENT_D = RGBColor(0x15, 0x31, 0x77)
ACCENT_SOFT = RGBColor(0xE7, 0xF2, 0xFF)
INK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x5A, 0x65, 0x73)
LINE = RGBColor(0xE3, 0xE6, 0xEA)
PLATE = RGBColor(0xF1, 0xF1, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x26, 0xAD, 0x50)
RED = RGBColor(0x66, 0x14, 0x29)
DIV_BG = RGBColor(0x15, 0x31, 0x77)
DIV_LIGHT = RGBColor(0x8F, 0xC4, 0xFF)
FONT = "Montserrat"

# Сетка слайда (16:9 шаблона)
MX = Cm(1.63)
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)
CONTENT_W = SLIDE_W - MX * 2

L_BLANK = None  # используем чистый layout «Текстовый» (без призраков)


def _layout(prs, name):
    for L in prs.slide_layouts:
        if L.name == name:
            return L
    return prs.slide_layouts[0]


def _del_slides(prs):
    lst = prs.slides._sldIdLst  # noqa: SLF001
    for sid in list(lst):
        prs.part.drop_rel(sid.attrib[qn('r:id')])
        lst.remove(sid)


def _strip(slide):
    """Убрать ВСЕ плейсхолдеры layout, чтобы не торчали подсказки."""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)  # noqa: SLF001


def _new(prs, layout_name="Текстовый"):
    s = prs.slides.add_slide(_layout(prs, layout_name))
    _strip(s)
    return s


def _txt(slide, x, y, w, h, text, *, size=12, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def _rect(slide, x, y, w, h, fill, *, line=None, shadow=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    if not shadow:
        sh.shadow.inherit = False
    return sh


def _roundrect(slide, x, y, w, h, fill, *, radius=0.07):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _send_back(slide, shape):
    tree = slide.shapes._spTree  # noqa: SLF001
    tree.remove(shape._element)
    tree.insert(2, shape._element)


# ── базовые компоненты «как в HTML» ──────────────────────────────

def _accent_bar(slide):
    """Левая вертикальная синяя полоса 10 px ширины — как в HTML."""
    bar = _rect(slide, Cm(0), Cm(0), Cm(0.3), SLIDE_H, ACCENT)
    _send_back(slide, bar)


def _kicker(slide, text, x=None, y=Cm(1.0)):
    _txt(slide, x or MX, y, CONTENT_W, Cm(0.6),
         text.upper(), size=11, bold=True, color=ACCENT)


def _h2(slide, text, y=Cm(1.8), size=28):
    _txt(slide, MX, y, CONTENT_W, Cm(1.6),
         text, size=size, bold=True, color=INK)


def _foot(slide, source):
    line = slide.shapes.add_connector(1, MX, SLIDE_H - Cm(1.5),
                                       SLIDE_W - MX, SLIDE_H - Cm(1.5))
    line.line.color.rgb = LINE; line.line.width = Pt(0.75)
    _txt(slide, MX, SLIDE_H - Cm(1.3), Cm(28), Cm(0.7),
         "Источник: " + source, size=10, color=GRAY)


def _card(slide, x, y, w, h, label, value, delta="", delta_color=GRAY,
          *, value_size=24):
    _roundrect(slide, x, y, w, h, PLATE, radius=0.08)
    pad = Cm(0.55)
    _txt(slide, x + pad, y + pad, w - 2 * pad, Cm(0.9),
         label, size=11, bold=True, color=GRAY)
    _txt(slide, x + pad, y + Cm(1.6), w - 2 * pad, h - Cm(2.8),
         value, size=value_size, bold=True, color=INK)
    if delta:
        _txt(slide, x + pad, y + h - pad - Cm(0.8), w - 2 * pad, Cm(0.7),
             delta, size=11, bold=True, color=delta_color)


def _cards_row(slide, y, items, h=Cm(4.2)):
    gap = Cm(0.4)
    n = len(items)
    w = (CONTENT_W - gap * (n - 1)) / n
    for i, it in enumerate(items):
        _card(slide, MX + (w + gap) * i, y, w, h, *it)
    return y + h


def _table(slide, x, y, w, headers, rows, *, fs=11, row_h=Cm(0.78),
           col_widths=None):
    cols = len(headers); n = len(rows) + 1
    t = slide.shapes.add_table(n, cols, x, y, w, row_h * n).table
    for r in t.rows:
        r.height = row_h
    # выключить встроенную тему таблицы (красим сами)
    tblPr = t._tbl.tblPr  # noqa: SLF001
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            t.columns[i].width = int(w * cw / total)

    def fill(cell, color):
        cell.fill.solid(); cell.fill.fore_color.rgb = color

    def cell_txt(cell, text, *, color=INK, bold=False, align=PP_ALIGN.LEFT):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Cm(0.2); cell.margin_right = Cm(0.2)
        cell.margin_top = Cm(0.05); cell.margin_bottom = Cm(0.05)
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.name = FONT; r.font.size = Pt(fs)
        r.font.bold = bold; r.font.color.rgb = color

    for c, h in enumerate(headers):
        cell = t.cell(0, c); fill(cell, ACCENT)
        cell_txt(cell, h, color=WHITE, bold=True,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
    for ri, row in enumerate(rows, 1):
        zebra = (ri % 2 == 0)
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            fill(cell, PLATE if zebra else WHITE)
            txt = str(val).strip()
            color = INK
            if "%" in txt and txt.startswith("+"): color = GREEN
            elif "%" in txt and (txt.startswith("−") or txt.startswith("-")):
                color = RED
            cell_txt(cell, val, color=color, bold=(ci == 0),
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
    return t


def _line_chart(slide, x, y, w, h, cats, series, *, title=None):
    cd = CategoryChartData(); cd.categories = cats
    for s in series:
        cd.add_series(s["name"], s["data"])
    ch = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, cd).chart
    ch.has_title = bool(title)
    if title:
        ch.chart_title.text_frame.text = title
        for r in ch.chart_title.text_frame.paragraphs[0].runs:
            r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = GRAY
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10); ch.legend.font.name = FONT
    ch.legend.font.color.rgb = GRAY
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = PLATE
    va.major_gridlines.format.line.width = Pt(0.5)
    va.format.line.fill.background()
    va.tick_labels.font.size = Pt(10); va.tick_labels.font.name = FONT
    va.tick_labels.font.color.rgb = GRAY
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.name = FONT
    ca.tick_labels.font.color.rgb = GRAY
    cols = [RGBColor(0xC2, 0xC8, 0xD0), ACCENT, ACCENT_D]
    for i, ser in enumerate(ch.series):
        ser.smooth = False
        col = cols[i] if i < len(cols) else ACCENT
        ser.format.line.color.rgb = col
        ser.format.line.width = Pt(2.5 if series[i].get("plan") else 3.25)
        if series[i].get("forecast"):
            ser.format.line.dash_style = 7
        m = ser.marker; m.style = 8; m.size = 7
        m.format.fill.solid(); m.format.fill.fore_color.rgb = col
        m.format.line.color.rgb = col
    return ch


# ── ТИПЫ СЛАЙДОВ ──────────────────────────────────────────────────

def title_slide(prs, title, kicker, sub, meta):
    s = _new(prs)
    _accent_bar(s)
    # «Логотип»
    dot = _roundrect(s, MX, Cm(2.0), Cm(1.0), Cm(1.0), ACCENT, radius=0.3)
    _txt(s, MX + Cm(1.3), Cm(2.05), Cm(8), Cm(1.0), "Контур",
         size=22, bold=True, color=INK)
    _txt(s, MX, Cm(4.5), CONTENT_W, Cm(0.8), kicker.upper(),
         size=12, bold=True, color=ACCENT)
    _txt(s, MX, Cm(5.4), CONTENT_W, Cm(4), title,
         size=60, bold=True, color=INK, line_spacing=1.05)
    _txt(s, MX, Cm(11.0), CONTENT_W, Cm(1.4),
         sub, size=14, color=GRAY)
    # 4 мета-блока
    for i, (k, v) in enumerate(meta):
        x = MX + Cm(i * 7.6)
        _txt(s, x, Cm(13.7), Cm(7.2), Cm(0.6), k,
             size=10, bold=True, color=GRAY)
        _txt(s, x, Cm(14.4), Cm(7.2), Cm(1.2), v,
             size=18, bold=True, color=INK)
    _foot(s, "Контур · продажи")
    return s


def divider(prs, num, title, sub):
    s = _new(prs)
    bg = _rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, DIV_BG)
    _send_back(s, bg)
    _txt(s, SLIDE_W - Cm(11), Cm(0.5), Cm(10), Cm(8),
         num, size=160, bold=True,
         color=RGBColor(0x24, 0x44, 0x86), align=PP_ALIGN.RIGHT)
    _txt(s, MX, Cm(6.2), CONTENT_W, Cm(0.8), "РАЗДЕЛ",
         size=12, bold=True, color=DIV_LIGHT)
    _txt(s, MX, Cm(7.2), Cm(26), Cm(5), title,
         size=54, bold=True, color=WHITE, line_spacing=1.05)
    _txt(s, MX, Cm(13.5), Cm(26), Cm(3), sub,
         size=14, color=RGBColor(0xC9, 0xDB, 0xF5))
    return s


def kpi_slide(prs, kicker, title, cards, *, source=""):
    s = _new(prs)
    _accent_bar(s)
    _kicker(s, kicker)
    _h2(s, title)
    _cards_row(s, Cm(5.0), cards)
    if source:
        _foot(s, source)
    return s


def table_slide(prs, kicker, title, headers, rows, *, source="",
                cards=None, fs=11, row_h=Cm(0.78), col_widths=None,
                table_y=None):
    s = _new(prs)
    _accent_bar(s)
    _kicker(s, kicker)
    _h2(s, title)
    y = Cm(5.0)
    if cards:
        y = _cards_row(s, y, cards) + Cm(0.6)
    if table_y:
        y = table_y
    _table(s, MX, y, CONTENT_W, headers, rows,
           fs=fs, row_h=row_h, col_widths=col_widths)
    if source:
        _foot(s, source)
    return s


def chart_slide(prs, kicker, title, cats, series, *,
                source="", cards=None):
    s = _new(prs)
    _accent_bar(s)
    _kicker(s, kicker)
    _h2(s, title)
    y = Cm(5.0)
    if cards:
        y = _cards_row(s, y, cards) + Cm(0.6)
    h = SLIDE_H - y - Cm(2.0)
    _line_chart(s, MX, y, CONTENT_W, h, cats, series)
    if source:
        _foot(s, source)
    return s


def bullets_slide(prs, kicker, title, bullets, *, source=""):
    s = _new(prs)
    _accent_bar(s)
    _kicker(s, kicker)
    _h2(s, title)
    y = Cm(5.2)
    for b in bullets:
        # «bullet» — синий квадратик слева
        _rect(s, MX, y + Cm(0.3), Cm(0.3), Cm(0.3), ACCENT)
        _txt(s, MX + Cm(0.7), y, CONTENT_W - Cm(0.7), Cm(2.0),
             b, size=15, color=INK, line_spacing=1.3)
        y += Cm(1.7)
    if source:
        _foot(s, source)
    return s


def two_col_bullets(prs, kicker, title, left_title, left, right_title, right,
                    *, source=""):
    s = _new(prs)
    _accent_bar(s)
    _kicker(s, kicker)
    _h2(s, title)
    cw = (CONTENT_W - Cm(1.0)) / 2
    # левая колонка
    _txt(s, MX, Cm(5.2), cw, Cm(0.7), left_title.upper(),
         size=12, bold=True, color=GREEN)
    y = Cm(6.0)
    for b in left:
        _rect(s, MX, y + Cm(0.3), Cm(0.3), Cm(0.3), GREEN)
        _txt(s, MX + Cm(0.7), y, cw - Cm(0.7), Cm(2.5),
             b, size=13, color=INK, line_spacing=1.3)
        y += Cm(2.2)
    # правая колонка
    rx = MX + cw + Cm(1.0)
    _txt(s, rx, Cm(5.2), cw, Cm(0.7), right_title.upper(),
         size=12, bold=True, color=RED)
    y = Cm(6.0)
    for b in right:
        _rect(s, rx, y + Cm(0.3), Cm(0.3), Cm(0.3), RED)
        _txt(s, rx + Cm(0.7), y, cw - Cm(0.7), Cm(2.5),
             b, size=13, color=INK, line_spacing=1.3)
        y += Cm(2.2)
    if source:
        _foot(s, source)
    return s


# ── СБОРКА — точное зеркало HTML-дека (26 слайдов) ───────────────

def build(template_path, out_path):
    prs = Presentation(template_path)
    _del_slides(prs)
    M = ["Янв", "Фев", "Мар", "Апр", "Май", "Июн",
         "Июл", "Авг", "Сен", "Окт", "Ноя", "Дек"]

    # 1. Титул
    title_slide(prs, "Итоги Q1 2026",
                "Квартальный отчёт по продажам",
                "Розница · Общепит · Кассы · ОФД · Маркет · Контекстная реклама · прогноз 2026",
                [("Период", "Q1 2026"),
                 ("Сравнение", "Q4 2025 · Q1 2025"),
                 ("Выручка факт", "68,2 млн ₽"),
                 ("Выполнение плана", "92%")])

    # 2. Обзор — карточки + таблица направлений
    kpi_slide(prs, "Обзор", "Ключевые метрики квартала",
              [("Выручка факт", "68,2 млн", "92% от плана 74,6 млн", RED),
               ("Кол-во оплат", "6 859", "план 6 773 · 101%", GRAY),
               ("Средний чек", "9,9 тыс", "факт / кол-во оплат", GRAY),
               ("Недобор к плану", "−6,3 млн", "в основном март", RED)],
              source="Царь свод (план-факт) + Царь продажи")

    # 3. Динамика
    table_slide(prs, "Помесячная динамика",
                "Бизнес-юнит: план и факт по месяцам, млн ₽",
                ["Месяц", "План", "Факт", "% выполнения", "Оплат"],
                [["Январь", "22,33", "21,12", "95%", "2 123"],
                 ["Февраль", "24,75", "23,69", "96%", "2 510"],
                 ["Март", "27,47", "23,42", "85%", "2 226"],
                 ["Итого Q1", "74,56", "68,23", "92%", "6 859"]],
                source="Царь свод все разрезы",
                table_y=Cm(5.0))

    # 4. Раздел 01
    divider(prs, "01", "Направления",
            "Розница · Общепит · Кассы · ОФД — выручка, динамика QoQ/YoY и ключевые тарифы.")

    # 5. План/факт по БЮ
    table_slide(prs, "Бизнес-юнит · план / факт",
                "Q1 2026 по направлениям · план 74,6 → факт 68,2 млн (92%)",
                ["Направление", "Факт Q1, млн ₽", "Доля", "Ключевая динамика"],
                [["Розница", "~32,5", "48%", "Оптим. Розница +9% QoQ; Премиум −42% YoY"],
                 ["Кассы", "~22,2", "33%", "Перерег. ККТ +48% YoY; Терминал +29% QoQ"],
                 ["Общепит", "~12,8", "19%", "Оптим. Общепит +44% QoQ — растёт"],
                 ["ОФД (срез)", "11,7", "—", "+4,2% QoQ, но −4,9% YoY"]],
                source="Царь свод + Царь продажи",
                col_widths=[3, 2, 1, 6], table_y=Cm(5.0))

    # 6–9. Направления
    for kicker, title, rows, src in [
        ("Розница", "Розница · ~32,5 млн ₽ (48% БЮ)",
         [["Оптимальный Розница", "7,13", "6,13", "6,66", "+9%", "−7%"],
          ["Все госсистемы", "3,97", "3,80", "3,47", "−9%", "−13%"],
          ["Базовый Розница", "1,89", "1,59", "1,41", "−11%", "−25%"],
          ["Премиум Розница", "1,64", "1,13", "0,95", "−16%", "−42%"],
          ["Маркировка", "1,39", "3,26", "1,86", "−43%", "+34%"]],
         "Царь продажи · БЮ = Госсистемы для розницы"),
        ("Общепит", "Общепит · ~12,8 млн ₽ (19% БЮ) · растёт",
         [["Оптимальный Общепит", "1,20", "0,89", "1,29", "+44%", "+8%"],
          ["КМ модификатор ЕГАИС", "1,17", "0,99", "1,09", "+10%", "−7%"],
          ["Базовый Общепит", "—", "—", "0,49", "·", "·"]],
         "Царь продажи · БЮ = Госсистемы для общепита"),
        ("Кассы", "Кассы · ~22,2 млн ₽ (33% БЮ) · самое здоровое",
         [["Терминал", "2,63", "2,14", "2,75", "+29%", "+5%"],
          ["Модификатор Маркировка", "2,46", "2,30", "2,55", "+11%", "+4%"],
          ["Лицензия Атол ИТС", "1,23", "1,39", "1,59", "+14%", "+29%"],
          ["Перерегистрация ККТ", "1,05", "1,30", "1,56", "+20%", "+48%"],
          ["Атол ИТС", "1,48", "1,00", "1,24", "+24%", "−17%"],
          ["Сканер", "1,21", "1,37", "1,16", "−15%", "−5%"]],
         "Царь продажи · БЮ = Кассовики"),
        ("ОФД", "ОФД · 11,69 млн ₽ (+4,2% QoQ · −4,9% YoY) · 2 122 оплаты",
         [["ОФД-36", "4,77", "4,79", "3,93", "−18%", "−18%"],
          ["ОФД-15", "4,17", "3,00", "3,73", "+24%", "−11%"],
          ["ОФД-13", "2,72", "2,73", "3,38", "+24%", "+24%"],
          ["ОФД-24 и прочие", "—", "—", "0,65", "·", "·"]],
         "Царь продажи · Тариф содержит «ОФД»"),
    ]:
        table_slide(prs, "Направление · " + kicker, title,
                    ["Тариф", "Q1 2025", "Q4 2025", "Q1 2026", "QoQ", "YoY"],
                    rows, source=src, table_y=Cm(5.0))

    # 10–12. Топ тарифов раздельно
    table_slide(prs, "Тарифы · Маркет", "Топ тарифов Маркета: QoQ и YoY",
                ["Тариф (Контур.Маркет)", "Q1 2025", "Q4 2025", "Q1 2026", "QoQ", "YoY"],
                [["Оптимальный Розница", "7,13", "6,13", "6,66", "+9%", "−7%"],
                 ["Все госсистемы", "3,97", "3,80", "3,47", "−9%", "−13%"],
                 ["Настройка Контур.Маркет", "1,59", "1,45", "2,93", "+102%", "+85%"],
                 ["Маркировка", "1,39", "3,26", "1,86", "−43%", "+34%"],
                 ["Базовый Розница", "1,89", "1,59", "1,41", "−11%", "−25%"],
                 ["Оптимальный Общепит", "1,20", "0,89", "1,29", "+44%", "+8%"],
                 ["КМ модификатор ЕГАИС", "1,17", "0,99", "1,09", "+10%", "−7%"],
                 ["Премиум Розница", "1,64", "1,13", "0,95", "−16%", "−42%"],
                 ["Базовый Услуги", "1,04", "0,67", "0,86", "+29%", "−18%"]],
                source="Царь продажи · продуктовые тарифы Контур.Маркет",
                fs=11, row_h=Cm(0.7), table_y=Cm(5.0))
    table_slide(prs, "Тарифы · ОФД", "Топ тарифов ОФД: QoQ и YoY",
                ["Тариф (Контур.ОФД)", "Q1 2025", "Q4 2025", "Q1 2026", "QoQ", "YoY"],
                [["ОФД-36", "4,77", "4,79", "3,93", "−18%", "−18%"],
                 ["ОФД-15", "4,17", "3,00", "3,73", "+24%", "−11%"],
                 ["ОФД-13", "2,72", "2,73", "3,38", "+24%", "+24%"],
                 ["ОФД-24", "—", "—", "0,46", "·", "·"],
                 ["Код активации ОФД-13", "—", "—", "0,10", "·", "·"],
                 ["ОФД-18", "—", "—", "0,07", "·", "·"],
                 ["Контур.ОФД-7 дней", "—", "—", "0,02", "·", "·"]],
                source="Царь продажи · Тариф содержит «ОФД»",
                table_y=Cm(5.0))
    table_slide(prs, "Тарифы · Кассы", "Топ тарифов Касс: QoQ и YoY",
                ["Тариф (Кассы)", "Q1 2025", "Q4 2025", "Q1 2026", "QoQ", "YoY"],
                [["Терминал", "2,63", "2,14", "2,75", "+29%", "+5%"],
                 ["Модификатор Маркировка", "2,46", "2,30", "2,55", "+11%", "+4%"],
                 ["Лицензия Атол ИТС", "1,23", "1,39", "1,59", "+14%", "+29%"],
                 ["Перерегистрация ККТ", "1,05", "1,30", "1,56", "+20%", "+48%"],
                 ["Атол ИТС", "1,48", "1,00", "1,24", "+24%", "−17%"],
                 ["Сертификат", "0,86", "0,66", "1,23", "+88%", "+43%"],
                 ["Сканер", "1,21", "1,37", "1,16", "−15%", "−5%"]],
                source="Царь продажи · БЮ = Кассовики",
                table_y=Cm(5.0))

    # 13. Факторы — две колонки
    two_col_bullets(prs, "Факторы квартала · почему так",
                    "Что двигало результат: успехи и сложности",
                    "Успехи",
                    ["Маркировочные волны: «Маркировка» +34% YoY, спрос толкает продажи.",
                     "Кросс-продуктовые запуски: «Настройка Контур.Маркет» +102% QoQ / +85% YoY.",
                     "Кассы — здоровый рост: Перерег. ККТ +48% YoY, Терминал +29% QoQ.",
                     "«тс ПИОТ»: [добавить суть инициативы и вклад]."],
                    "Сложности",
                    ["Март провалил план (85%): рост плана до 27,5 млн, факт на уровне февраля.",
                     "Маркет −21% YoY; Премиум/Базовый Розница −42% / −25% YoY.",
                     "ОФД −4,9% YoY; реклама: ROMI ≈ −50%, CPO Маркета 46,8 тыс ₽."],
                    source="Царь продажи + аудит рекламы")

    # 14. Раздел 02
    divider(prs, "02", "Прогноз 2026",
            "Маркет и ОФД. Ансамбль моделей · сценарии низ/база/верх.")

    # 15. Маркет KPI
    kpi_slide(prs, "Прогноз · выручка", "Маркет — выполнение плана 2026",
              [("Годовой план", "119,0 млн", "Продукт = Маркет", GRAY),
               ("Факт за 5 мес", "44,5 млн", "91% плана периода", RED),
               ("YoY к 2025", "−21%", "падение спроса", RED),
               ("Прогноз года", "91%", "108,3 млн · коридор 86–96%", RED)],
              source="план-факт «Продукт = Маркет» · ансамбль моделей")

    # 16. Маркет график
    chart_slide(prs, "Прогноз · выручка",
                "Маркет — план / факт / прогноз 2026, млн ₽", M,
                [{"name": "План", "plan": True, "data":
                  [9.78, 9.46, 10.98, 10.45, 8.36, 10.32, 10.13, 9.01, 8.79, 9.95, 9.71, 12.10]},
                 {"name": "Факт", "data":
                  [8.24, 9.20, 10.26, 8.99, 7.78, None, None, None, None, None, None, None]},
                 {"name": "Прогноз", "forecast": True, "data":
                  [None, None, None, None, 7.78, 9.41, 9.24, 8.22, 8.02, 9.08, 8.86, 11.04]}],
                source="ансамбль · pacing темпа YTD = 91%")

    # 17. ОФД KPI
    kpi_slide(prs, "Прогноз · выручка", "ОФД — выполнение плана 2026",
              [("Годовой план", "37,5 млн", "проект п453", GRAY),
               ("Факт за 5 мес", "18,9 млн", "123% плана периода", GREEN),
               ("YoY к 2025", "−16%", "план консервативен", GRAY),
               ("Прогноз года", "115%", "43,3 млн · коридор 103–124%", GREEN)],
              source="план-факт проект п453 · ансамбль моделей")

    # 18. ОФД график
    chart_slide(prs, "Прогноз · выручка",
                "ОФД — план / факт / прогноз 2026, млн ₽", M,
                [{"name": "План", "plan": True, "data":
                  [2.81, 3.11, 3.18, 3.24, 2.97, 3.04, 3.03, 2.99, 2.90, 3.28, 3.20, 3.79]},
                 {"name": "Факт", "data":
                  [4.05, 4.12, 4.18, 3.71, 2.84, None, None, None, None, None, None, None]},
                 {"name": "Прогноз", "forecast": True, "data":
                  [None, None, None, None, 2.84, 3.6, 3.4, 3.1, 3.1, 3.4, 3.2, 4.5]}],
                source="план-факт проект п453 · ансамбль моделей")

    # 19. Прогноз лидов
    kpi_slide(prs, "Прогноз · лиды", "Прогноз по лидам 2026 (мой расчёт)",
              [("Маркет — лиды", "~16,1 тыс", "нужно ~17,8 тыс · дефицит ~10%", RED),
               ("CR2 лид→оплата", "36% / 46%", "Маркет / ОФД", GRAY),
               ("ОФД — лиды", "~15,4 тыс", "минимум ~15,0 тыс · в плане", GREEN),
               ("Оплаты прогноз", "5,8 / 7,1 тыс", "Маркет / ОФД", GRAY)],
              source="мой прогноз = оплаты ÷ CR2")

    # 20. Раздел 03
    divider(prs, "03", "Контекстная реклама",
            "Расход, CPL, CPO и привлечённые оплаты по Маркету и ОФД. Дашборды кампаний, янв 2025 – май 2026.")

    # 21. Маркет vs ОФД
    table_slide(prs, "Эффективность контекста",
                "Контекст: Маркет vs ОФД (янв 2025 – май 2026)",
                ["Показатель", "Маркет", "ОФД", "Разница"],
                [["Расход", "14,7 млн ₽", "7,1 млн ₽", "Маркет ×2,1"],
                 ["Оплаты привлечено", "315", "397", "ОФД +26%"],
                 ["CPO — стоимость оплаты", "46,8 тыс ₽", "17,9 тыс ₽", "ОФД ×2,6 дешевле"],
                 ["CPL — средний", "10,0 тыс ₽", "4,5 тыс ₽", "ОФД ×2,2 дешевле"],
                 ["Выручка по целевым лидам", "3,9 млн ₽", "2,8 млн ₽", "—"],
                 ["ДРР (расход / выручка)", "≈ 3,8×", "≈ 2,5×", "ОФД эффективнее"]],
                source="дашборд Контекст · суммы сверены с «Итого»",
                row_h=Cm(1.0), fs=12, table_y=Cm(5.0))

    # 22. Контекст по продуктам
    table_slide(prs, "Контекст · сводка",
                "Привлечение по продуктам — оплаты по месяцам",
                ["Месяц", "ОФД", "Маркет", "Бандл"],
                [["М1", "54", "63", "52"],
                 ["М2", "65", "95", "46"],
                 ["М3", "50", "80", "60"],
                 ["М4", "52", "36", "43"],
                 ["М5", "37", "28", "68"],
                 ["М6", "8", "13", "1"],
                 ["Итого", "266", "315", "270"]],
                source="data (15) ОФД · data (16) Маркет · data (17) Бандл",
                table_y=Cm(5.0))

    # 23. Раздел 04
    divider(prs, "04", "Маркетинг: каналы и планы 2026",
            "Контекст 2023–2026, брендформанс, SEO и аудит платной рекламы.")

    # 24. Годовая динамика
    table_slide(prs, "Каналы · контекстная реклама",
                "Контекст: годовая динамика 2023–2026 (РК Маркет + бандл)",
                ["Показатель", "2023", "2024", "2025", "2026 план"],
                [["Расход, млн ₽", "32,6", "32,6", "24,7", "17,0"],
                 ["Лиды", "3 986", "3 572", "3 205", "2 900"],
                 ["Оплаты", "536", "527", "538", "500"],
                 ["CR в оплату", "13,4%", "14,8%", "16,8%", "17,2%"],
                 ["Средний чек, тыс ₽", "33,6", "35,7", "36,7", "38,0"],
                 ["Выручка, млн ₽", "18,0", "18,8", "19,7", "19,0"],
                 ["Выручка − расход, млн ₽", "−14,5", "−13,7", "−5,0", "+2,0"]],
                source="план по каналам, Порубов",
                table_y=Cm(5.0))

    # 25. Брендформанс + SEO
    kpi_slide(prs, "Каналы · брендформанс и SEO",
              "Брендформанс и SEO — планы 2026",
              [("Брендформанс · факт", ">2 млн ₽", "за окт–дек 2025", GRAY),
               ("Брендформанс · план 2026", "~11 млн ₽", "из них ~6,5 — Маркет", GREEN),
               ("SEO · 2025", "+3%", "при норме +20%", RED),
               ("SEO · план 2026", "~+10 млн ₽", "при росте со 2 квартала", GREEN)],
              source="план по каналам, Порубов")

    # 26. Аудит Петровой
    kpi_slide(prs, "Аудит платной рекламы · Петрова",
              "Аудит контекста: вклад в выручку и ROMI",
              [("Отключение всей рекламы", "−12%", "выручки · ~12 млн ₽/год", RED),
               ("Вклад Бандла в ОФД", "32%", "выручки ОФД из рекламы", GRAY),
               ("ROMI · реалистично", "~−50%", "в плюс выйти потенциала нет", RED),
               ("ROMI в мае", "+15%", "оптимизация РК + продажи", GREEN)],
              source="«Маркет+ОФД Аудит полный», М. Петрова")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


if __name__ == "__main__":
    import sys
    t = sys.argv[1] if len(sys.argv) > 1 else \
        "Шаблон презентации Контур Blue 2023_16x9_Montserrat (2).pptx"
    o = sys.argv[2] if len(sys.argv) > 2 else "Итоги Q1 2026 — B · как в HTML.pptx"
    build(t, o)
    print(f"OK -> {o}")
